import streamlit as st
from io import BytesIO
from fpdf import FPDF
from PIL import Image
from docx import Document
from pptx import Presentation
import openpyxl
import os


def convert_to_pdf(file):
    try:
        file_type = file.type
        file_name = file.name

        if file_type == "application/pdf":
            # If already a PDF, return the original content
            return file.getvalue(), file_name

        elif file_type == "text/plain":
            # Convert text files to PDF
            pdf = FPDF()
            pdf.add_page()
            pdf.set_font("Arial", size=12)
            content = file.read().decode("utf-8")
            for line in content.splitlines():
                pdf.cell(200, 10, txt=line, ln=True)
            pdf_bytes = BytesIO()
            pdf.output(pdf_bytes)
            return pdf_bytes.getvalue(), file_name.replace(".txt", ".pdf")

        elif file_type == "image/png" or file_type == "image/jpeg":
            # Convert image files (JPG/PNG) to PDF
            image = Image.open(file)
            pdf_bytes = BytesIO()
            image.save(pdf_bytes, format="PDF")
            return pdf_bytes.getvalue(), file_name.replace(".png", ".pdf").replace(".jpg", ".pdf")

        elif file_type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
            # Convert Word files to PDF
            doc = Document(file)
            pdf = FPDF()
            pdf.set_auto_page_break(auto=True, margin=15)
            pdf.add_page()
            pdf.set_font("Arial", size=12)
            for para in doc.paragraphs:
                pdf.multi_cell(0, 10, para.text)
            pdf_bytes = BytesIO()
            pdf.output(pdf_bytes)
            return pdf_bytes.getvalue(), file_name.replace(".docx", ".pdf")

        elif file_type == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet":
            # Convert Excel files to PDF
            wb = openpyxl.load_workbook(file, data_only=True)
            sheet = wb.active
            pdf = FPDF()
            pdf.add_page()
            # pdf.set_font("Arial", size=12)
            pdf.set_font(family='Times',style='B',size=12)
            for row in sheet.iter_rows(values_only=True):
                row_text = "\t".join([str(cell) if cell is not None else "" for cell in row])
                pdf.multi_cell(0, 10, row_text)
            pdf_bytes = BytesIO()
            pdf.output(pdf_bytes)
            return pdf_bytes.getvalue(), file_name.replace(".xlsx", ".pdf")

        elif file_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
            # Convert PowerPoint files to PDF
            ppt = Presentation(file)
            pdf = FPDF()
            pdf.set_auto_page_break(auto=True, margin=15)
            for slide in ppt.slides:
                pdf.add_page()
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            pdf.multi_cell(0, 10, paragraph.text)
            pdf_bytes = BytesIO()
            pdf.output(pdf_bytes)
            return pdf_bytes.getvalue(), file_name.replace(".pptx", ".pdf")

        else:
            return None, None

    except Exception as e:
        st.error(f"Error converting {file.name}: {e}")
        return None, None


st.title("File to PDF Converter")

# File uploader
uploaded_files = st.file_uploader(
    "Upload your files (Word, Excel, PowerPoint, Image, Text, or PDF)", 
    accept_multiple_files=True
)

if uploaded_files:
    st.header("Converted Files")
    for uploaded_file in uploaded_files:
        print(type(uploaded_file))
        pdf_data, pdf_name = convert_to_pdf(uploaded_file)
        if pdf_data:
            st.download_button(
                label=f"Download {pdf_name}",
                data=pdf_data,
                file_name=pdf_name,
                mime="application/pdf"
            )
        else:
            st.warning(f"Could not convert {uploaded_file.name}. Unsupported file type.")
